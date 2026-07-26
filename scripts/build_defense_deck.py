"""Build the Shishoza capstone-defense deck.

Every figure is an existing project export (unchanged), every number is taken from
the saved metrics JSON/CSV files. Writes Shishoza_Defense_Deck.pptx in the repo root.
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "Shishoza_Defense_Deck.pptx"

# ------------------------------------------------------------------ design tokens
W, H = 20.0, 11.25                      # slide size, inches
M = 1.15                                # page margin
CW = W - 2 * M                          # content width

CREAM = RGBColor(0xF6, 0xF3, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x45, 0x5E, 0x28)
GREEN_DK = RGBColor(0x2C, 0x3E, 0x19)
GREEN_TINT = RGBColor(0xEA, 0xF0, 0xE1)
INK = RGBColor(0x12, 0x13, 0x0F)
MUTED = RGBColor(0x6B, 0x6B, 0x60)
LINE = RGBColor(0xDD, 0xD8, 0xC9)
AMBER = RGBColor(0xB2, 0x6B, 0x18)
RISK = RGBColor(0xC0, 0x49, 0x2F)

FONT = "Avenir Next"
FOOTER = "Shishoza  ·  Capstone Defense  ·  Izere Uwonkunda Marie Claire Kerie"

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------- helpers
def new_slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=True, adj=0.06):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        shape.adjustments[0] = adj
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    return shape


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         spacing=1.15, anchor=MSO_ANCHOR.TOP, space_after=0, italic=False,
         caps_track=False):
    """runs: a string, or a list of (text, {overrides}) tuples, or a list of strings."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        opts = {}
        if isinstance(item, tuple):
            item, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("spacing", spacing)
        p.space_after = Pt(opts.get("space_after", space_after))
        p.space_before = Pt(opts.get("space_before", 0))
        r = p.add_run()
        r.text = item
        f = r.font
        f.name = opts.get("font", FONT)
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", italic)
        f.color.rgb = opts.get("color", color)
    return tb


def track(tb, spc=180):
    """Letter-space every run in a textbox (spc is in 1/100 pt)."""
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font._rPr.set("spc", str(spc))
    return tb


def est_lines(txt, width_in, size, ratio=0.55):
    """Estimated wrapped line count for a body font at `size` pt in `width_in`."""
    cpl = max(12, int(width_in * 72 / (size * ratio)))
    n = 0
    for para in txt.split("\n"):
        n += max(1, -(-len(para) // cpl))
    return n


def eyebrow(s, label, color=GREEN, y=0.78):
    track(text(s, M, y, CW, 0.35, label.upper(), size=12.5, color=color, bold=True))


def heading(s, title, sub=None, y=1.28, color=INK, rule=True, size=36, subw=0.86):
    text(s, M, y, CW, 1.0, title, size=size, color=color, bold=True, spacing=1.05)
    yy = y + est_lines(title, CW, size, 0.53) * (size * 1.16 / 72.0) + 0.20
    if rule:
        rect(s, M, yy, 1.9, 0.055, fill=GREEN if color == INK else WHITE, radius=False)
        yy += 0.34
    if sub:
        text(s, M, yy, CW * subw, 0.6, sub, size=17, color=MUTED, spacing=1.25)
        yy += est_lines(sub, CW * subw, 17) * (17 * 1.25 / 72.0) + 0.30
    else:
        yy += 0.18
    return yy


def footer(s, n, color=MUTED):
    text(s, M, H - 0.72, CW * 0.7, 0.3, FOOTER, size=10.5, color=color)
    text(s, W - M - 1.2, H - 0.72, 1.2, 0.3, str(n), size=10.5, color=color,
         align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, items, size=17.5, gap=13, dot=GREEN, lead=0.0):
    """items: list of str, or (str, 'sub') for a secondary line."""
    yy = y
    for it in items:
        kind = "main"
        if isinstance(it, tuple):
            it, kind = it
        if kind == "head":
            track(text(s, x, yy, w, 0.34, it.upper(), size=12.5, color=GREEN, bold=True))
            yy += 0.52
            continue
        sz = size if kind == "main" else size - 2
        rect(s, x + 0.03, yy + 0.115, 0.13, 0.13,
             fill=dot if kind == "main" else LINE, radius=True, adj=0.5)
        text(s, x + 0.42, yy - 0.04, w - 0.42, 0.4, it, size=sz,
             color=INK if kind == "main" else MUTED, spacing=1.22)
        yy += est_lines(it, w - 0.42, sz) * (sz * 1.22 / 72.0) + gap / 72.0
    return yy


def picture(s, path, x, y, w, h, card=True, pad=0.22, card_fill=WHITE, valign="center"):
    """Fit an image inside the box, centred, optionally on a card."""
    iw, ih = Image.open(ROOT / path).size
    ar = iw / ih
    bw, bh = (w - 2 * pad, h - 2 * pad) if card else (w, h)
    if bw / bh > ar:
        ph = bh
        pw = ph * ar
    else:
        pw = bw
        ph = pw / ar
    px = x + (w - pw) / 2
    py = y + (0 if valign == "top" else (h - ph) / 2) + (pad if valign == "top" and card else 0)
    if card:
        rect(s, px - pad, py - pad, pw + 2 * pad, ph + 2 * pad,
             fill=card_fill, line=LINE, lw=0.75, adj=0.03)
    s.shapes.add_picture(str(ROOT / path), Inches(px), Inches(py),
                         Inches(pw), Inches(ph))
    return px, py, pw, ph


def takeaway(s, msg, y=None, h=1.18, tone="green"):
    y = H - 1.05 - h if y is None else y
    fill, col = (GREEN_TINT, GREEN_DK) if tone == "green" else (RGBColor(0xFA, 0xEE, 0xE2), AMBER)
    rect(s, M, y, CW, h, fill=fill, radius=True, adj=0.12)
    rect(s, M, y, 0.09, h, fill=GREEN if tone == "green" else AMBER, radius=False)
    text(s, M + 0.42, y, CW - 0.9, h, msg, size=16.5, color=col,
         anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)


def stat_tiles(s, y, tiles, h=2.05, size=44):
    """tiles: list of (value, label)."""
    n = len(tiles)
    gap = 0.35
    w = (CW - gap * (n - 1)) / n
    for i, (val, lab) in enumerate(tiles):
        x = M + i * (w + gap)
        rect(s, x, y, w, h, fill=WHITE, line=LINE, lw=0.75, adj=0.08)
        text(s, x + 0.45, y + 0.34, w - 0.9, 0.9, val, size=size, color=GREEN,
             bold=True, spacing=1.0)
        text(s, x + 0.45, y + 0.34 + size / 72.0 + 0.16, w - 0.9, 0.8, lab,
             size=14, color=MUTED, spacing=1.2)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


N = 0


def page(bg=CREAM, numbered=True):
    global N
    N += 1
    s = new_slide(bg)
    if numbered:
        footer(s, N, RGBColor(0xA9, 0xBE, 0x91) if bg == GREEN else MUTED)
    return s


# ========================================================================= slides
# 1 — title
s = new_slide(GREEN)
N = 1
rect(s, 0, 0, W, 0.28, fill=GREEN_DK, radius=False)
track(text(s, M, 2.55, CW, 0.4, "AFRICAN LEADERSHIP UNIVERSITY",
     size=13.5, color=RGBColor(0xC7, 0xD6, 0xB4), bold=True), 260)
text(s, M, 3.15, CW * 0.94, 1.4, "Shishoza", size=88, color=WHITE, bold=True, spacing=0.95)
text(s, M, 4.75, CW * 0.86, 1.8,
     "A machine-learning system for predicting deforestation risk\nat the land-parcel level in Rwanda",
     size=30, color=RGBColor(0xE7, 0xEE, 0xDC), spacing=1.22)
rect(s, M, 6.95, 2.2, 0.06, fill=WHITE, radius=False)
text(s, M, 7.45, 8.0, 2.0, [
    ("Izere Uwonkunda Marie Claire Kerie", {"size": 22, "bold": True, "color": WHITE}),
    ("BSc Software Engineering  ·  Capstone Defense", {"size": 16, "color": RGBColor(0xC7, 0xD6, 0xB4), "space_before": 6}),
], spacing=1.2)
text(s, W - M - 6.5, 7.45, 6.5, 2.0, [
    ("Supervisor", {"size": 13, "color": RGBColor(0xC7, 0xD6, 0xB4), "bold": True, "align": PP_ALIGN.RIGHT}),
    ("Ndinelao Iitumba", {"size": 22, "bold": True, "color": WHITE, "align": PP_ALIGN.RIGHT, "space_before": 4}),
], spacing=1.2)
text(s, M, H - 0.95, CW, 0.35, "kerie1-shishoza.hf.space   ·   github.com/izerekerie/shishoza_ml_project",
     size=13, color=RGBColor(0xA9, 0xBE, 0x91))
notes(s, "Open with the one-line framing: Rwanda's 2024 Forest Law requires a permit "
         "to cut trees at any size, but no tool can see clearing at the size Rwandan "
         "parcels actually are. Shishoza is that screening tool.")

# 2 — problem context
s = page()
eyebrow(s, "Background & problem context")
y = heading(s, "A new legal duty at any size — and no tool that can see it",
            "Rwanda has grown national forest cover. What remains is small, scattered clearing: "
            "smallholders taking part of a plot.")
stat_tiles(s, y + 0.05, [
    ("0.18 ha", "average registered farm plot in Rwanda\n(Mugabowindekwe et al., 2024)"),
    ("0.1–6.25 ha", "minimum mapping unit of global alert systems\n(Ygorra et al., 2024)"),
    ("Any size", "permit now required under the 2024 Forest Law\n(Law No. 046/2024, replacing the 2 ha threshold)"),
])
y2 = y + 2.45
bullets(s, M, y2, CW * 0.55, [
    ("The two gaps stack", "head"),
    "LEGAL gap — closed. The old 2013 law only required a permit at 2 ha and above; the 2024 law removed the threshold.",
    "DETECTION gap — open. Global systems (GFW, RADD, DETER) work well at the scale they were built for, but sub-hectare clearing stays invisible to them.",
    "Rwanda's 0.18 ha parcel falls below the detection floor, so officers must decide permits on tiny parcels with nothing to screen them.",
])
bullets(s, M + CW * 0.6, y2, CW * 0.4, [
    ("Stakeholders", "head"),
    "Smallholder farmers and citizens applying to cut",
    "Forest officers and district managers deciding permits",
    "The Rwanda Forestry Authority",
])
notes(s, "Lead with the gap: a new legal duty exists at any size, but no tool can see the "
         "sub-hectare clearing it applies to. Be careful — the global tools are NOT bad at "
         "large scale, they work well there. The point is they miss the SMALL clearings that "
         "dominate in Rwanda.")

# 3 — literature
s = page()
eyebrow(s, "Literature review & relevant works")
y = heading(s, "What exists, and what is missing for Rwanda")
bullets(s, M, y + 0.1, CW * 0.52, [
    ("What exists", "head"),
    "Amoakoh et al. (2021) — multi-source remote sensing + Random Forest, Ghana peatland.",
    "Gutkin et al. (2023) — supervised tree-cover classification in Eastern Rwanda, 85–92% forest-class accuracy.",
    "Ygorra et al. (2024) — near-real-time tropical deforestation alerts (CuSum); the F1 = 0.71 baseline this work is measured against.",
    "Jocea et al. (2025) — review of 89 Sentinel-2 land-cover studies; models lose 15–25% accuracy out of region.",
])
x2 = M + CW * 0.56
rect(s, x2 - 0.4, y + 0.02, CW * 0.44 + 0.5, 4.55, fill=GREEN_TINT, radius=True, adj=0.05)
bullets(s, x2, y + 0.3, CW * 0.44, [
    ("The gaps this project addresses", "head"),
    "No published model we could identify is trained on Rwanda's own satellite data for deforestation detection.",
    "No test of Sentinel-2 + Sentinel-1 + SRTM fusion in Rwanda's fragmented montane landscape.",
    "No citizen-facing, parcel-level tool tied to the permit process.",
])
takeaway(s, "The out-of-region accuracy loss reported by Jocea et al. is exactly why this model is "
            "trained on Rwandan data and evaluated with spatial cross-validation rather than a random split.")
notes(s, "Keep this short. The point of the slide is the third gap: nobody has tied a "
         "parcel-level model to the permit process.")

# 4 — research questions
s = page()
eyebrow(s, "Research questions")
y = heading(s, "Four questions the system had to answer")
qs = [
    ("RQ1", "How effectively can multi-source satellite fusion detect sub-hectare deforestation at the parcel level in Rwanda?", "Answered — F1 0.832 with all three sources"),
    ("RQ2", "How does detection accuracy hold under spatial cross-validation, on regions the model never saw?", "Answered — F1 0.753 ± 0.081"),
    ("RQ3", "What value does a 500 m neighbourhood analysis add to parcel-level risk?", "Answered — 12% of parcels escalated"),
    ("RQ4", "Can the system deliver a usable parcel-level risk signal to citizens and officers?", "Answered — deployed, p95 = 270 ms"),
]
yy = y + 0.15
for tag, q, a in qs:
    rect(s, M, yy, CW, 1.30, fill=WHITE, line=LINE, lw=0.75, adj=0.11)
    rect(s, M, yy, 0.10, 1.30, fill=GREEN, radius=False)
    text(s, M + 0.55, yy, 1.4, 1.30, tag, size=21, color=GREEN, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, M + 2.1, yy, CW - 6.6, 1.30, q, size=17.5, color=INK, spacing=1.2,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, W - M - 4.5, yy, 4.0, 1.30, a, size=14, color=GREEN_DK, bold=True,
         align=PP_ALIGN.RIGHT, spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
    yy += 1.50
notes(s, "Say the answers out loud here — the panel then knows where the talk is going, "
         "and every later slide is evidence for one of these four lines.")

# 5 — objectives & scope
s = page()
eyebrow(s, "Objectives & scope")
y = heading(s, "What the project set out to build — and what it deliberately excludes")
bullets(s, M, y + 0.1, CW * 0.46, [
    ("Objectives", "head"),
    "Build a Random Forest model fusing Sentinel-2 optical, Sentinel-1 radar and SRTM terrain data.",
    "Evaluate it honestly with spatial cross-validation, not a random split.",
    "Deliver a web application that returns a parcel-level risk output to a citizen and to an officer.",
])
x2 = M + CW * 0.52
bullets(s, x2, y + 0.1, CW * 0.46, [
    ("In scope", "head"),
    "Parcel-level clearing detection and HIGH / MEDIUM / LOW risk tiering",
    "500 m neighbourhood context analysis",
    "Web app with land-certificate or parcel-shape upload",
    "District-scoped officer review of citizen requests",
])
bullets(s, x2, y + 3.5, CW * 0.46, [
    ("Out of scope", "head"),
    ("Environmental impact assessment — the system detects clearing and assigns risk, it does not judge impact", "sub"),
    ("GPS positioning — location comes from the uploaded certificate or shape", "sub"),
    ("Direct land-registry integration", "sub"),
    ("Real-time continuous monitoring", "sub"),
])
notes(s, "The included/excluded split protects you: it pre-empts 'why doesn't it do X' questions.")

# 6 — proposed solution
s = page()
eyebrow(s, "Proposed solution")
y = heading(s, "Shishoza — parcel-level risk screening, in three steps",
            "A citizen uploads a land certificate or a screenshot of the parcel shape. No GPS, no land-registry access, no account.")
steps = [
    ("1", "Locate the parcel", "The certificate or shape image is read and the parcel location and area are extracted."),
    ("2", "Predict clearing", "A Random Forest trained on fused optical + radar + terrain data returns a clearing probability."),
    ("3", "Tier the risk", "A 500 m neighbourhood check adds spatial context, and the parcel is labelled HIGH, MEDIUM or LOW."),
]
sw = (CW - 0.7) / 3
for i, (n, t, d) in enumerate(steps):
    x = M + i * (sw + 0.35)
    rect(s, x, y + 0.1, sw, 2.75, fill=WHITE, line=LINE, lw=0.75, adj=0.07)
    rect(s, x + 0.5, y + 0.5, 0.75, 0.75, fill=GREEN, radius=True, adj=0.5)
    text(s, x + 0.5, y + 0.63, 0.75, 0.5, n, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + 0.5, y + 1.5, sw - 1.0, 0.5, t, size=21, color=INK, bold=True)
    text(s, x + 0.5, y + 2.05, sw - 1.0, 1.0, d, size=14.5, color=MUTED, spacing=1.25)
tiers = [("HIGH", RISK), ("MEDIUM", AMBER), ("LOW", GREEN)]
yy = y + 3.3
text(s, M, yy, 4.0, 0.4, "RISK TIERS RETURNED", size=12.5, color=GREEN, bold=True)
for i, (lab, col) in enumerate(tiers):
    x = M + i * 2.6
    rect(s, x, yy + 0.55, 2.2, 0.65, fill=col, radius=True, adj=0.18)
    text(s, x, yy + 0.55, 2.2, 0.65, lab, size=16, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
takeaway(s, "Why machine learning: sub-hectare, cloud-prone, fragmented clearing is exactly where "
            "fixed-threshold tools fail. A model trained on fused optical + radar + terrain data is the appropriate approach.")
notes(s, "Be precise: the system detects clearing and assigns risk — it does not assess "
         "environmental impact.")

# 7 — requirements
s = page()
eyebrow(s, "System requirements")
y = heading(s, "Functional and non-functional requirements",
            "Gathered from document analysis of the 2013 and 2024 forest laws and the literature, plus a stakeholder interview.")
bullets(s, M, y + 0.05, CW * 0.46, [
    ("Functional", "head"),
    "Upload a land certificate or parcel-shape image and extract the parcel location.",
    "Run the model and return a HIGH / MEDIUM / LOW risk result for that parcel.",
    "Provide 500 m neighbourhood risk context around the parcel.",
    "Let a citizen submit a cutting request; let an officer review and decide it.",
    "Authenticated forest-manager access scoped to their own district.",
])
x2 = M + CW * 0.52
bullets(s, x2, y + 0.05, CW * 0.46, [
    ("Non-functional", "head"),
    "Privacy — citizens use the app anonymously: no account, SHA-256-hashed IPs, no identity stored. Officers log in, scoped to their district.",
    "Performance — a parcel result returned within a few seconds.",
    "Usability — a citizen reaches a result with no account and minimal steps.",
    "Portability — web-based, with an open documented API for institutional integration.",
])
notes(s, "Requirements gathering: document analysis (laws, literature) plus the stakeholder "
         "interview. If asked about sample size, be straight — it was one interview plus "
         "document analysis, and the acceptance test had five users.")

# 8 — architecture
s = page()
eyebrow(s, "System design")
text(s, M, 1.26, CW, 0.7, "System architecture", size=36, color=INK, bold=True)
picture(s, "results/diagrams/architecture.png", M, 2.15, CW, 8.20, card=False)
notes(s, "Walk it in the order a request travels: certificate upload in the browser, Flask API, "
         "extractor and geometry resolver, model loader, risk classifier, neighbourhood check, "
         "result. Then point at the bottom band and say the model itself was trained offline in "
         "notebooks and published to Hugging Face, which is why the app only loads an artefact.")

# 9 — data model
s = page()
eyebrow(s, "System design")
text(s, M, 1.26, CW, 0.7, "Data model — live SQLite schema", size=36, color=INK, bold=True)
picture(s, "results/diagrams/erd.png", M, 2.15, CW, 7.9, card=False, valign="top")
notes(s, "Six entities. Say explicitly that predictions_log.csv is a flat monitoring file "
         "and the sector boundaries are GeoJSON — they are not tables. The report's earlier "
         "entity count was reconciled to this live schema.")

# 9b — classes
s = page()
eyebrow(s, "System design")
text(s, M, 1.26, CW, 0.7, "Application classes", size=36, color=INK, bold=True)
picture(s, "results/diagrams/class_diagram.png", M, 2.15, CW, 7.9, card=False, valign="top")
notes(s, "Thirteen classes. The boundary class (FlaskAPI) is the only entry point; control "
         "classes coordinate, services do the work, entities hold state. The review, "
         "notification and logging classes are the five added after the report draft.")

# 10 — implementation
s = page()
eyebrow(s, "Implementation")
y = heading(s, "Built as a Flask application around a trained Random Forest")
cols = [
    ("Model & data", ["Random Forest, 400 trees (scikit-learn)",
                      "Google Earth Engine data pipeline",
                      "23,319 nationally sampled pixels",
                      "17 features from three sources",
                      "Trained offline in Jupyter notebooks"]),
    ("Web application", ["Flask + Jinja2 templates",
                         "HTML / CSS / JavaScript front end",
                         "Leaflet for parcel and sector maps",
                         "Flasgger for the public API docs",
                         "Anonymous citizen flow, no account"]),
    ("Storage", ["SQLite — accounts and citizen requests",
                 "Analysis cache and simulation runs",
                 "predictions_log.csv for drift monitoring",
                 "GeoJSON for 416 sector boundaries",
                 "IPs stored only as SHA-256 hashes"]),
    ("Deployment", ["Docker + gunicorn",
                    "Free 16 GB Hugging Face Space",
                    "kerie1-shishoza.hf.space",
                    "Model pulled from the HF Hub at build",
                    "Public REST API documented in Swagger"]),
]
cwid = (CW - 3 * 0.35) / 4
for i, (t, items) in enumerate(cols):
    x = M + i * (cwid + 0.35)
    rect(s, x, y + 0.05, cwid, 4.55, fill=WHITE, line=LINE, lw=0.75, adj=0.06)
    text(s, x + 0.45, y + 0.45, cwid - 0.9, 0.5, t, size=19, color=GREEN, bold=True)
    yy = y + 1.25
    for it in items:
        rect(s, x + 0.45, yy + 0.10, 0.10, 0.10, fill=LINE, radius=True, adj=0.5)
        text(s, x + 0.75, yy - 0.03, cwid - 1.2, 0.5, it, size=14.5, color=INK, spacing=1.2)
        yy += est_lines(it, cwid - 1.2, 14.5) * (14.5 * 1.22 / 72.0) + 0.16
takeaway(s, "The 325 MB model is too large for GitHub, so the Docker image pulls it from a public "
            "Hugging Face model repository at build time. Citizens stay anonymous — the IP is stored only as an irreversible SHA-256 hash.")
notes(s, "If you show a live demo here, warm the Space first so it responds quickly.")

# 11 — data sourcing & EDA
s = page()
eyebrow(s, "Data sourcing & exploratory analysis")
y = heading(s, "23,319 pixels sampled across the whole country",
            "Sentinel-2 optical, Sentinel-1 radar and SRTM terrain via Google Earth Engine · labels from Hansen Global Forest Change (loss 2020–2024).",
            rule=False)
picture(s, "results/eda/01_class_balance.png", M, y + 0.05, CW * 0.46, 5.9, valign="top")
picture(s, "results/eda/05_spatial_distribution.png", M + CW * 0.52, y + 0.05, CW * 0.48, 5.9, valign="top")
takeaway(s, "Near-balanced by design — 12,003 cleared against 11,316 stable pixels, 51.5% positive — so the reported F1 "
            "is not inflated by class imbalance. The sample is national, not one forest, which is what makes the spatial cross-validation meaningful.")
notes(s, "Lead with the class balance: it pre-empts the 'is your F1 inflated by imbalance?' "
         "question before it is asked. Then say the sample is national, because that is what "
         "makes the spatial cross-validation on the results slides meaningful.")

# 12 — EDA features
s = page()
eyebrow(s, "Exploratory analysis")
y = heading(s, "The features separate the two classes — but not cleanly", rule=False)
picture(s, "results/eda/02_feature_distributions.png", M, y + 0.05, CW * 0.42, 7.1, valign="top")
picture(s, "results/eda/03_correlation_heatmap.png", M + CW * 0.45, y + 0.05, CW * 0.34, 7.1, valign="top")
bullets(s, M + CW * 0.81, y + 0.45, CW * 0.19, [
    ("What this shows", "head"),
    "Cleared and stable pixels overlap in every single feature.",
    "No threshold on NDVI, NBR or backscatter separates them on its own.",
    "That overlap is the argument for a trained model rather than a fixed rule.",
])
notes(s, "This is the slide that justifies machine learning instead of an index threshold. "
         "The distributions overlap; a single cut-off on NDVI would misclassify a large "
         "share of pixels either way.")

# 13 — headline result
s = page()
eyebrow(s, "Results — RQ2")
y = heading(s, "The honest number is 0.753, not 0.836",
            "Spatial block cross-validation: ten K-means blocks held out by GroupKFold, so whole regions are unseen at test time.",
            rule=False)
picture(s, "results/metrics/nat_chart5_spatial_cv.png", M, y - 0.05, CW * 0.63, 5.85, valign="top")
picture(s, "results/metrics/nat_chart7_confusion_matrix.png", M + CW * 0.65, y - 0.05, CW * 0.35, 5.85, valign="top")
takeaway(s, "0.832 on a random split drops to 0.753 ± 0.081 when whole regions are held out — and 0.753 still beats "
            "the 0.71 global baseline of Ygorra et al. (2024) under a stricter evaluation than most published work reports.")
notes(s, "Lead with 0.753, never 0.836. Do not mix the Nyungwe pair (0.788 random / 0.733 "
         "spatial) with the national pair (0.832 / 0.753) — quoting one number from each is "
         "the easiest mistake to make and the easiest for a panelist to catch. The fold spread "
         "0.617 to 0.838 is real and you should volunteer it.")

# 14 — where it generalises
s = page()
eyebrow(s, "Results — generalisation")
y = heading(s, "Where it holds, and where it does not", rule=False)
picture(s, "results/metrics/nat_chart6_province_lopo.png", M, y + 0.05, CW * 0.66, 5.9, valign="top")
bullets(s, M + CW * 0.70, y + 0.35, CW * 0.30, [
    ("Leave-one-province-out", "head"),
    "Kigali City 0.850, North 0.806, East 0.805, South 0.806.",
    "West drops to 0.709 with recall 0.646 — the weakest region and the honest limit of the model.",
    "The West is the steep, cloud-prone, montane edge of the country; it is also where the fold spread in spatial CV comes from.",
])
takeaway(s, "Four of five provinces sit at or above 0.80. The West is where I would target the next round of "
            "training data — and it is the reason the spatial-CV standard deviation is ± 0.081 rather than ± 0.01.")
notes(s, "Volunteering the weakest province is stronger than being asked for it. If pressed on "
         "why the West is hard: steeper terrain, more persistent cloud, and more fragmented "
         "smallholder mosaics.")

# 15 — RQ1 fusion
s = page()
eyebrow(s, "Results — RQ1")
y = heading(s, "Fusing three data sources beats any one of them",
            "Identical folds, identical model, four feature configurations.", rule=False)
picture(s, "results/metrics/nat_chart2_rq1_experiments.png", M, y + 0.05, CW * 0.60, 5.7, valign="top")
picture(s, "results/metrics/nat_chart4_importance_by_source.png", M + CW * 0.63, y + 0.05, CW * 0.37, 5.7, valign="top")
takeaway(s, "Optical alone 0.755 → plus terrain 0.821 → plus radar 0.794 → all three 0.832. Terrain adds more than radar, "
            "which is the counter-intuitive finding: clearing concentrates at accessible elevations and gentler slopes. Elevation is the single most important feature.")
notes(s, "The counter-intuitive result is the interesting one — say it plainly. Terrain is not "
         "detecting the clearing, it is telling the model where clearing is plausible. Optical "
         "51%, terrain 25%, radar 24% of total importance.")

# 16 — RQ2 patch size
s = page()
eyebrow(s, "Results — RQ2")
y = heading(s, "It detects the parcel sizes global systems miss",
            "Recall against real Hansen clearing patches, sized with connectedPixelCount; out-of-fold predictions.",
            rule=False)
picture(s, "results/metrics/nat_chart9_rq2_patchsize.png", M, y + 0.05, CW * 0.64, 5.8, valign="top")
bullets(s, M + CW * 0.68, y + 0.3, CW * 0.32, [
    ("Recall by real patch size", "head"),
    "0.771 at 0.09–0.18 ha (n = 1,023) — the size global alert systems cannot see.",
    "Rises monotonically to 0.873 above 1.8 ha.",
    "Overall recall 0.838 across 11,543 real patches.",
    "Target for the sub-hectare band was 0.75 — the dashed line on the chart.",
])
takeaway(s, "This chart is the project's core claim: 0.771 recall at 0.09–0.18 ha, the band that falls below the "
            "minimum mapping unit of every global alert system, and above the 0.75 target set for it.")
notes(s, "This is the answer to 'so what?'. The whole argument of the dissertation is that "
         "small clearing is invisible to existing tools, and this is the measurement that "
         "shows the model sees it.")

# 17 — RQ3 neighbourhood
s = page()
eyebrow(s, "Results — RQ3")
y = heading(s, "Neighbourhood context changes the decision for 1 parcel in 8",
            "3,000 sampled national parcels · 25-nearest-neighbour parcel probability plus a 500 m surrounding check.",
            rule=False)
picture(s, "results/metrics/rq3_neighbourhood.png", M, y + 0.05, CW * 0.60, 5.8, valign="top")
bullets(s, M + CW * 0.64, y + 0.35, CW * 0.36, [
    ("What the 500 m check adds", "head"),
    "Parcel-only view: 1,586 LOW, 1,414 HIGH — nothing in between.",
    "With neighbourhood context: 1,226 LOW, 360 MEDIUM, 1,414 HIGH.",
    "360 of 3,000 parcels (12.0%) were raised from LOW to MEDIUM.",
    "Those are cumulative-pressure cases: a quiet parcel inside an actively clearing 500 m radius.",
])
takeaway(s, "A parcel-only permit check would have passed all 360 of those parcels as LOW risk. Context is what turns a "
            "per-parcel probability into a decision an officer can act on.")
notes(s, "Be precise about the direction: the neighbourhood only escalates here, it never "
         "de-escalates. That is a deliberate design choice for a screening tool.")

# 18 — model selection
s = page()
eyebrow(s, "Model selection")
y = heading(s, "Why Random Forest is the deployed model",
            "Three families, identical 17 features, identical folds. Grid search over 96 combinations first moved the Random Forest F1 by 0.004 — within noise.",
            rule=False)
picture(s, "results/experiments/model_family_comparison.png", M, y + 0.05, CW * 0.62, 5.6, valign="top")
bullets(s, M + CW * 0.66, y + 0.15, CW * 0.34, [
    ("Spatial-CV F1 — the error bars overlap", "head"),
    "XGBoost 0.767 ± 0.068  ·  LightGBM 0.757 ± 0.072  ·  Random Forest 0.752 ± 0.078",
    "A 1.5-point spread against a fold-to-fold standard deviation of ~0.07: the differences sit inside the noise.",
    "Model size: Random Forest 32.75 MB, LightGBM 1.22 MB, XGBoost 0.58 MB.",
])
takeaway(s, "The ceiling is the Hansen labels — noisy at sub-hectare scale and spatially autocorrelated — not model capacity. "
            "Random Forest stayed in production because every reported result was produced with it; adopting XGBoost is the first item of future work.",
         tone="amber")
notes(s, "Volunteer this slide — do not wait to be asked. Never say you chose Random Forest "
         "because it is smaller or faster; it is the largest and slowest of the three and your "
         "own notebook says so. Concede that XGBoost scored highest, then win on the reasoning: "
         "the gap is inside the noise, and the validated model is the defensible one to keep.")

# 19 — external check
s = page()
eyebrow(s, "External check — stated honestly")
y = heading(s, "Agreement with Hansen across 365 sectors is weak, and I report it as weak",
            rule=False)
picture(s, "results/metrics/nat_chart10_hansen_validation.png", M, y + 0.05, CW * 0.58, 5.8, valign="top")
bullets(s, M + CW * 0.62, y + 0.3, CW * 0.38, [
    ("Pearson r = 0.32  ·  Spearman ρ = 0.29", "head"),
    "This is a consistency check, not independent validation — Hansen also supplies the training labels.",
    "The disagreements concentrate where the model flags clearing below Hansen's minimum mapping unit.",
    "That is consistent with the gap this project targets, but it is not proof of it without field data.",
    "Field ground-truth collection is the stated route to closing this, and it is in future work.",
])
takeaway(s, "Volunteering a weak correlation and explaining why it is weak is stronger than defending a number that "
            "cannot carry the weight. Independent validation needs field data this project did not have.", tone="amber")
notes(s, "If asked 'how strong is your external validation?': r = 0.32 across 365 sectors, "
         "weak, and I report it as weak. Two reasons — Hansen supplies my training labels, so "
         "it is a consistency check; and the disagreements sit exactly where I would expect, "
         "below Hansen's minimum mapping unit.")

# 20 — testing & evaluation
s = page()
eyebrow(s, "Testing & evaluation")
y = heading(s, "Measured against targets set before the build",
            "Unit tests on the data pipeline and model inference · integration tests on end-to-end parcel lookup · system tests on the full upload-to-result workflow.",
            rule=False)
rows = [
    ("Metric", "Target", "Achieved", "Status"),
    ("Spatial-CV F1", "≥ 0.71  (Ygorra et al. 2024 baseline)", "0.753", "Passed"),
    ("Sub-hectare recall, 0.09–0.18 ha", "≥ 0.75", "0.771", "Passed"),
    ("Manager risk API response", "< 2 s", "p95 = 270 ms at 200 concurrent users, 0.38% failures", "Passed"),
    ("Citizen live-satellite analysis", "< 2 s", "~18 s warm, ~96 s cold", "Not met"),
    ("Map payload on a weak link", "usable on 3G", "168 KB, ~3.4 s on Slow 3G", "Passed"),
    ("Usability acceptance test", "qualitative", "5 participants — indicative, not representative", "Indicative"),
]
colx = [0.0, 5.6, 8.9, 15.3]
colw = [5.4, 3.1, 6.4, 1.5]
yy = y + 0.05
for i, r in enumerate(rows):
    head = i == 0
    rh = 0.72 if head else 0.78
    if not head:
        rect(s, M, yy, CW, rh, fill=WHITE if i % 2 else RGBColor(0xFB, 0xFA, 0xF5),
             line=LINE, lw=0.6, adj=0.03)
    for j, cell in enumerate(r):
        col = GREEN if head else INK
        if not head and j == 3:
            col = RISK if cell == "Not met" else (AMBER if cell == "Indicative" else GREEN)
        text(s, M + 0.4 + colx[j], yy, colw[j], rh, cell,
             size=13 if head else 15, color=col,
             bold=head or j == 3, spacing=1.15, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.RIGHT if j == 3 else PP_ALIGN.LEFT)
    yy += rh + (0.10 if head else 0.06)
takeaway(s, "The live Earth Engine path is too slow for the citizen flow and I list it as not met. The shipped default is the "
            "nearest-sample fallback, which returns sub-second; caching and region-restricting the composite is the fix, and it is in future work.",
         tone="amber", y=yy + 0.12, h=1.0)
notes(s, "Include the row you did not pass. A table where every row says Passed reads as "
         "unmeasured; one honest failure with a stated fix reads as engineering. Load-test "
         "evidence is in results/performance/RQ4_delivery.md.")

# 21 — demo
s = page(bg=GREEN)
track(text(s, M, 1.1, CW, 0.4, "LIVE DEMO", size=13.5, color=RGBColor(0xC7, 0xD6, 0xB4), bold=True))
text(s, M, 1.75, CW * 0.8, 1.2, "The system, running", size=46, color=WHITE, bold=True)
links = [
    ("Citizen flow", "kerie1-shishoza.hf.space/citizen"),
    ("Manager review", "kerie1-shishoza.hf.space/login"),
    ("API documentation", "kerie1-shishoza.hf.space/apidocs"),
]
for i, (t, u) in enumerate(links):
    x = M + i * ((CW - 0.7) / 3 + 0.35)
    ww = (CW - 0.7) / 3
    rect(s, x, 3.5, ww, 1.6, fill=GREEN_DK, radius=True, adj=0.10)
    text(s, x + 0.5, 3.82, ww - 1.0, 0.5, t, size=20, color=WHITE, bold=True)
    text(s, x + 0.5, 4.38, ww - 1.0, 0.5, u, size=14, color=RGBColor(0xC7, 0xD6, 0xB4))
text(s, M, 5.7, CW * 0.6, 0.4, "WHAT I WILL SHOW", size=12.5, color=RGBColor(0xC7, 0xD6, 0xB4), bold=True)
flow = ["Upload a land certificate or parcel shape",
        "Location extracted from the document — no GPS",
        "Analysis runs and returns HIGH / MEDIUM / LOW",
        "Submit the request to the district officer",
        "Officer reviews and decides; the citizen is emailed"]
yy = 6.3
for i, f in enumerate(flow):
    rect(s, M + 0.03, yy + 0.13, 0.13, 0.13, fill=RGBColor(0xC7, 0xD6, 0xB4), radius=True, adj=0.5)
    text(s, M + 0.45, yy - 0.02, CW * 0.6, 0.4, f, size=17, color=WHITE)
    yy += 0.62
picture(s, "results/eda/swagger_ui.png", M + CW * 0.62, 5.6, CW * 0.38, 4.35, card=True)
notes(s, "Warm the Space before you present — it is a free tier and it sleeps. Demo the "
         "nearest-sample path, not the live Earth Engine pull: the live pull takes 18–96 s "
         "and will stall in front of the panel.")

# 22 — discussion
s = page()
eyebrow(s, "Results & discussion")
y = heading(s, "What was achieved, and what the results cannot claim",
            "Main outcome: a working parcel-level risk-screening system trained on Rwanda's own data, addressing a gap the 2024 permit law creates.")
bullets(s, M, y + 0.05, CW * 0.45, [
    ("Strengths", "head"),
    "Honest spatial-CV evaluation, reported instead of the flattering random-split figure.",
    "Sub-hectare detection at the scale global alert systems miss.",
    "Neighbourhood context that changes the tier for 12% of parcels.",
    "Privacy by design — anonymous for citizens, authenticated and district-scoped for officers.",
])
x2 = M + CW * 0.51
bullets(s, x2, y + 0.05, CW * 0.49, [
    ("Limitations", "head"),
    "Balanced training means the reported F1 is not operational precision — in the field clearing is rare, so precision will be lower than 0.75.",
    "Label dependence: training labels and the external cross-check both come from Hansen, so r = 0.32 is a consistency check, not independent proof.",
    "Accuracy is label-limited, not model-limited — three model families and a 96-combination grid search all land within noise.",
    "The live-satellite citizen path is too slow (18–96 s); the fast nearest-sample path is the shipped default.",
    "No direct land-registry access; cold-start latency on the free host; the USSD channel is designed but not deployed.",
])
notes(s, "Naming limitations honestly is a strength in a defense, not a weakness. If a panelist "
         "raises one you have already listed, you have won that exchange.")

# 23 — impact
s = page()
eyebrow(s, "Impact & significance")
y = heading(s, "Why this matters beyond the model")
items = [
    ("For smallholders", "Free, parcel-level satellite evidence about their own land, with no account and no smartphone GPS required."),
    ("For forest officers", "A screening signal they currently do not have, reducing blind manual review of permit applications."),
    ("For the 2024 Forest Law", "Support for the permit obligation at the scale it actually applies — any size, including 0.18 ha plots."),
    ("For the region", "A locally-trained, replicable method for comparable smallholder landscapes in East Africa."),
]
cwid = (CW - 3 * 0.35) / 4
for i, (t, d) in enumerate(items):
    x = M + i * (cwid + 0.35)
    rect(s, x, y + 0.1, cwid, 3.6, fill=WHITE, line=LINE, lw=0.75, adj=0.07)
    rect(s, x + 0.5, y + 0.6, 0.55, 0.09, fill=GREEN, radius=False)
    text(s, x + 0.5, y + 1.0, cwid - 1.0, 0.9, t, size=20, color=GREEN, bold=True, spacing=1.1)
    text(s, x + 0.5, y + 2.0, cwid - 1.0, 1.5, d, size=14.5, color=INK, spacing=1.3)
takeaway(s, "The relevant comparison is not perfection — it is the status quo of no tool at all at this scale, and a published "
            "global baseline of 0.71 that this work exceeds under a stricter evaluation.")
notes(s, "Keep this slide short and confident. It is the answer to 'so what', stated at the "
         "level of the country rather than the model.")

# 24 — conclusion & future work
s = page()
eyebrow(s, "Conclusion & future work")
y = heading(s, "Where this ends, and what comes next")
bullets(s, M, y + 0.05, CW * 0.45, [
    ("Key achievements", "head"),
    "Built and validated a Random Forest model with honest spatial-CV reporting, and confirmed the choice against XGBoost and LightGBM.",
    "Delivered a citizen-facing web application with a parcel-level risk output, deployed and load-tested.",
    "Added a 500 m neighbourhood analysis that measurably improves risk triage.",
    ("Lesson learned: match every claim to the evidence, report the conservative number, and keep the human in the decision.", "sub"),
])
x2 = M + CW * 0.51
bullets(s, x2, y + 0.05, CW * 0.49, [
    ("Future work, in priority order", "head"),
    "Adopt XGBoost in production and re-run the downstream analyses against it — same accuracy within noise, 0.58 MB instead of 325 MB.",
    "Collect field ground-truth so accuracy can be validated independently of Hansen.",
    "Cache and region-restrict the Earth Engine composite to bring the live citizen path under the worker timeout.",
    "Deploy the USSD channel for feature-phone inclusion.",
    "Pursue National Land Authority data access, and re-validate live-period accuracy against future ground-truth labels.",
])
notes(s, "End on the XGBoost item — it shows you know what you would do next and why, and it "
         "turns the model-choice question into a plan instead of a weakness.")

# 25 — thank you
s = new_slide(GREEN)
N += 1
rect(s, 0, 0, W, 0.28, fill=GREEN_DK, radius=False)
text(s, M, 3.2, CW, 1.3, "Thank you", size=76, color=WHITE, bold=True)
text(s, M, 4.75, CW * 0.6, 0.6, "Questions and answers", size=28, color=RGBColor(0xC7, 0xD6, 0xB4))
rect(s, M, 5.85, 2.2, 0.06, fill=WHITE, radius=False)
text(s, M, 6.5, CW * 0.55, 2.2, [
    ("Demo   kerie1-shishoza.hf.space/citizen", {"size": 19, "color": WHITE, "space_after": 10}),
    ("Repository   github.com/izerekerie/shishoza_ml_project", {"size": 19, "color": WHITE, "space_after": 10}),
    ("Supervisor   Ndinelao Iitumba", {"size": 19, "color": RGBColor(0xC7, 0xD6, 0xB4)}),
], spacing=1.2)
nums = [("0.753", "spatial-CV F1"), ("0.771", "recall at 0.09–0.18 ha"), ("12%", "parcels escalated"), ("270 ms", "p95 API response")]
for i, (v, l) in enumerate(nums):
    x = M + CW * 0.55 + i * 2.15
    text(s, x, 6.5, 2.0, 0.7, v, size=30, color=WHITE, bold=True)
    text(s, x, 7.15, 2.0, 0.8, l, size=12.5, color=RGBColor(0xC7, 0xD6, 0xB4), spacing=1.2)
text(s, M, H - 0.95, CW, 0.35, FOOTER, size=12, color=RGBColor(0xA9, 0xBE, 0x91))
notes(s, "Prepared answers: (1) XGBoost scored higher but inside the noise, and every "
         "downstream result was produced with the Random Forest. (2) The boosters were trained "
         "to test whether the ceiling was the model or the data — it is the labels. (3) External "
         "validation is weak at r = 0.32 and I say so. (4) 0.75 F1 is a screening signal, not a "
         "verdict; the officer still decides, and the alternative today is no tool at all.")

prs.save(OUT)
print(f"wrote {OUT.name} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
