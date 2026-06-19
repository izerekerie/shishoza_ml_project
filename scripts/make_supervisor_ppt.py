"""Builds TreeSight_Supervisor_Update.pptx — a short, decision-led deck for the
supervisor meeting. Content is drawn from the REVISED proposal + DECISIONS.md so it
matches what Kerie has actually built (Nyungwe, Hansen labels, 2020-22 train / 23-24 test)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- palette -------------------------------------------------------------
GREEN = RGBColor(0x1B, 0x5E, 0x20)   # forest green
DARK  = RGBColor(0x21, 0x21, 0x21)
GREY  = RGBColor(0x55, 0x55, 0x55)
AMBER = RGBColor(0xB7, 0x6E, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF8, 0xE9)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)  # rectangle bg
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tf


def setp(p, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, space=6, italic=False):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space)
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Calibri"


def bullet(tf, text, size=18, color=DARK, bold=False, level=0, space=8, italic=False):
    p = tf.add_paragraph(); p.level = level
    setp(p, text, size, color, bold, space=space, italic=italic)
    return p


def accent_bar(slide):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.25), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background(); bar.shadow.inherit = False


def title_block(slide, kicker, title):
    accent_bar(slide)
    tf = textbox(slide, 0.7, 0.45, 12.0, 1.5)
    setp(tf.paragraphs[0], kicker, 14, GREEN, bold=True, space=2)
    bullet(tf, title, 30, DARK, bold=True, space=4)


# ============================================================ SLIDE 1 — TITLE
s = add_slide(GREEN)
tf = textbox(s, 0.9, 2.2, 11.5, 3.2)
setp(tf.paragraphs[0], "TreeSight  ·  Supervisor Update", 18, LIGHT, bold=True, space=6)
bullet(tf, "Detecting small-scale forest clearing in", 34, WHITE, bold=True, space=0)
bullet(tf, "Rwanda's Nyungwe buffer zone", 34, WHITE, bold=True, space=18)
bullet(tf, "What I have built so far, my four research questions,", 18, LIGHT, space=2)
bullet(tf, "and the one decision I want your guidance on.", 18, LIGHT, space=2)
tf2 = textbox(s, 0.9, 6.5, 11.5, 0.6)
setp(tf2.paragraphs[0], "Kerie  ·  BSc Software Engineering, ALU  ·  Final-year capstone", 13, LIGHT)

# ====================================================== SLIDE 2 — FOUNDATION
s = add_slide()
title_block(s, "WHERE I AM", "What I have already decided and built")
tf = textbox(s, 0.7, 1.9, 12.0, 5.2)
setp(tf.paragraphs[0], "I am not starting from zero — the core pipeline and key decisions are in place:", 18, GREY, space=12)
for b in [
    "Model:  Random Forest classifier (chosen over deep learning — works on a laptop "
    "and gives feature-importance, which directly answers RQ1)",
    "Data:  Sentinel-2 optical + Sentinel-1 radar + SRTM terrain, from Google Earth Engine",
    "Labels (ground truth):  Hansen Global Forest Change — already chosen",
    "Area:  Nyungwe National Park buffer zone (the most-documented zone, with a published baseline)",
    "Time:  2020-2024, already split honestly — train on 2020-2022, test on 2023-2024",
    "App:  working web tool — upload land certificate -> parcel risk + recovery + cut simulation",
]:
    bullet(tf, b, 18, DARK, space=10)
tf2 = textbox(s, 0.7, 6.55, 12.0, 0.7)
setp(tf2.paragraphs[0],
     "Message: the foundation is solid. Today is about pressure-testing my reasoning, not rescuing it.",
     16, GREEN, bold=True, italic=True)

# ====================================================== SLIDE 3 — THE 4 RQs
s = add_slide()
title_block(s, "SAME PAGE", "My four research questions")
tf = textbox(s, 0.7, 1.9, 12.0, 5.3)
rqs = [
    ("RQ1 — Does radar help?",
     "Effect of adding Sentinel-1 radar to Sentinel-2 optical on the F1-score for detecting "
     "small-scale clearing in Nyungwe. (Tested via 4 data combinations.)"),
    ("RQ2 — How small can it see?",
     "Relationship between clearing patch size and detection accuracy, as size approaches "
     "Rwanda's typical farm of 0.18 ha (~18 pixels at 10 m)."),
    ("RQ3 — Recovery evidence",
     "What cumulative forest-recovery evidence the neighbourhood analysis adds beyond Rwanda's "
     "current permit-approval process."),
    ("RQ4 — Delivery to people",
     "How effectively the system delivers tree-cover loss + risk to citizens and forest managers "
     "at their GPS land parcel, before a permit application."),
]
for head, body in rqs:
    bullet(tf, head, 19, GREEN, bold=True, space=2)
    bullet(tf, body, 16, DARK, space=12)

# ====================================================== SLIDE 4 — WHY LABELS
s = add_slide()
title_block(s, "A QUESTION I ASKED MYSELF", "Why does the model need labels?")
tf = textbox(s, 0.7, 1.9, 7.4, 5.2)
setp(tf.paragraphs[0], "A label is the 'answer on the back of the flashcard.'", 19, DARK, bold=True, space=12)
for b in [
    "I show the model a 10 m tile -> it guesses 'cleared / not cleared'.",
    "The label tells it the true answer -> it corrects itself.",
    "Repeat over thousands of tiles -> it learns what clearing looks like.",
    "This is 'supervised' learning = learning from labelled answers.",
]:
    bullet(tf, b, 17, DARK, space=10)
bullet(tf, "Without labels the model sees pixels but cannot know which are clearings — "
           "there is nothing to learn from.", 17, GREY, space=12, italic=True)
# side panel
panel = s.shapes.add_shape(1, Inches(8.5), Inches(2.0), Inches(4.0), Inches(3.6))
panel.fill.solid(); panel.fill.fore_color.rgb = LIGHT; panel.line.color.rgb = GREEN
panel.line.width = Pt(1.5); panel.shadow.inherit = False
ptf = panel.text_frame; ptf.word_wrap = True
ptf.margin_left = Inches(0.25); ptf.margin_top = Inches(0.2)
setp(ptf.paragraphs[0], "My label source", 16, GREEN, bold=True, space=8)
bullet(ptf, "Hansen Global Forest Change", 18, DARK, bold=True, space=6)
bullet(ptf, "Free, peer-reviewed, global yearly forest-loss layer — already my ground truth. "
            "I do not hand-label tiles.", 15, GREY, space=4)

# ============================================== SLIDE 5 — CURRENT TRAINING SETUP
s = add_slide()
title_block(s, "WHAT I TRAIN WITH NOW", "My current training setup, in plain terms")
tf = textbox(s, 0.7, 1.9, 12.0, 5.3)
rows = [
    ("Input each pixel sees", "10 m Sentinel-2 colour + Sentinel-1 radar + SRTM elevation"),
    ("Answer key (labels)", "Hansen Global Forest Change loss/no-loss"),
    ("Where", "Nyungwe buffer zone — steep, cloudy, fragmented, 0.18 ha farms"),
    ("When", "2020-2024 (covers both rainy seasons, so seasonal change is learned)"),
    ("Honest test", "Train 2020-2022, test on unseen 2023-2024  (temporal hold-out)"),
    ("Experiments", "4 data combinations to isolate the radar effect (RQ1)"),
    ("App validation", "10 land parcels across different tree-loss levels"),
]
for k, v in rows:
    p = tf.add_paragraph(); p.space_after = Pt(11)
    r1 = p.add_run(); r1.text = k + ":  "
    r1.font.bold = True; r1.font.size = Pt(18); r1.font.color.rgb = GREEN; r1.font.name = "Calibri"
    r2 = p.add_run(); r2.text = v
    r2.font.size = Pt(18); r2.font.color.rgb = DARK; r2.font.name = "Calibri"

# ============================================== SLIDE 6 — THE ONE TENSION
s = add_slide()
title_block(s, "THE ONE OPEN ISSUE", "A wording tension I spotted in my own proposal")
tf = textbox(s, 0.7, 2.0, 12.0, 4.8)
bullet(tf, "My main objective says the app serves 'citizens' at 'any registered Rwanda land parcel.'",
       19, DARK, space=8)
bullet(tf, "But my scope trains and validates only on the Nyungwe buffer zone.",
       19, DARK, space=14)
bullet(tf, "So a citizen in Musanze or Kayonza would get a prediction from a model that never "
           "saw their landscape — and my own literature review (Shrestha 2025) says models lose "
           "15-25% accuracy outside their training region.", 18, AMBER, bold=True, space=14)
bullet(tf, "I caught this myself — it is a consistency issue between two paragraphs, not a flaw "
           "in the method.", 17, GREY, italic=True, space=4)

# ============================================== SLIDE 7 — OPTIONS + REC
s = add_slide()
title_block(s, "MY RECOMMENDATION", "Two ways to resolve it")
# Option A box
a = s.shapes.add_shape(1, Inches(0.7), Inches(2.0), Inches(5.7), Inches(4.4))
a.fill.solid(); a.fill.fore_color.rgb = LIGHT; a.line.color.rgb = GREEN
a.line.width = Pt(2); a.shadow.inherit = False
atf = a.text_frame; atf.word_wrap = True; atf.margin_left = Inches(0.3); atf.margin_top = Inches(0.25)
setp(atf.paragraphs[0], "Option A  (I recommend)", 18, GREEN, bold=True, space=4)
bullet(atf, "Narrow the app's wording to the Nyungwe buffer zone.", 17, DARK, bold=True, space=8)
for b in ["Keeps my tight, well-justified scope", "Fits the 10-week timeline",
          "Honest: I only claim where I validated", "Change = two paragraphs of text"]:
    bullet(atf, "+ " + b, 15, DARK, space=6)
# Option B box
b_ = s.shapes.add_shape(1, Inches(6.9), Inches(2.0), Inches(5.7), Inches(4.4))
b_.fill.solid(); b_.fill.fore_color.rgb = RGBColor(0xFB, 0xF3, 0xE0); b_.line.color.rgb = AMBER
b_.line.width = Pt(2); b_.shadow.inherit = False
btf = b_.text_frame; btf.word_wrap = True; btf.margin_left = Inches(0.3); btf.margin_top = Inches(0.25)
setp(btf.paragraphs[0], "Option B", 18, AMBER, bold=True, space=4)
bullet(btf, "Train nationally (sample a few zones per province).", 17, DARK, bold=True, space=8)
for b in ["App could honestly claim 'anywhere'", "- More data collection + labelling work",
          "- Risky inside a 10-week project", "- Dilutes the Nyungwe case study + baseline"]:
    col = DARK if b.startswith("App") else GREY
    bullet(btf, ("+ " if b.startswith("App") else "") + b, 15, col, space=6)

# ============================================== SLIDE 8 — GUIDANCE ASKS
s = add_slide(GREEN)
accent_bar(s)
tf = textbox(s, 0.7, 0.6, 12.0, 1.0)
setp(tf.paragraphs[0], "WHERE I WANT YOUR GUIDANCE", 14, LIGHT, bold=True, space=2)
bullet(tf, "Decisions I've made — please pressure-test them", 28, WHITE, bold=True)
tf2 = textbox(s, 0.7, 1.9, 12.0, 5.2)
asks = [
    "Scope: I lean toward narrowing the app to Nyungwe (Option A). Do you agree, or do you "
    "want the national route?",
    "RQ1: I isolate the radar effect with 4 data combinations + feature importance. Is that "
    "design convincing?",
    "RQ2: I report accuracy by patch size down to 0.18 ha. Headline metric — or do you want a "
    "precision/recall breakdown by size?",
    "Evaluation: I report F1 with a temporal hold-out (train 2020-22 / test 23-24). Is that "
    "the rigour you expect?",
    "Validation: 10 parcels within the zone — is that a defensible sample for a 10-week project?",
]
for a in asks:
    p = tf2.add_paragraph(); p.space_after = Pt(13)
    r = p.add_run(); r.text = "•  " + a
    r.font.size = Pt(18); r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf3 = textbox(s, 0.7, 6.7, 12.0, 0.6)
setp(tf3.paragraphs[0],
     "Opener: \"I've made some decisions and want to pressure-test them with you — "
     "there's one, scope, I'd most like your steer on.\"", 14, LIGHT, italic=True)

out = "TreeSight_Supervisor_Update.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
